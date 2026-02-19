import io
from typing import Final

import pdfplumber
from pptx import Presentation
from docx import Document
from fastapi import UploadFile

SUPPORTED_EXTENSIONS: Final[set[str]] = {".pdf", ".docx", ".pptx", ".txt"}


async def load_document(file: UploadFile) -> str:
    filename = (file.filename or "").lower()
    content = await file.read()

    if filename.endswith(".pdf"):
        return _load_pdf(content)
    if filename.endswith(".docx"):
        return _load_docx(content)
    if filename.endswith(".pptx"):
        return _load_pptx(content)
    if filename.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")

    raise ValueError(f"Unsupported file type. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}.")


async def load_documents(files: list[UploadFile]) -> str:
    text_parts: list[str] = []
    for file in files:
        doc_text = await load_document(file)
        if doc_text.strip():
            label = file.filename or "Document"
            text_parts.append(f"--- {label} ---\n{doc_text}")
    combined = "\n\n".join(text_parts).strip()
    if not combined:
        raise ValueError("All uploaded documents are empty after extraction.")
    return combined


def _load_pdf(content: bytes) -> str:
    text_parts: list[str] = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)
    return "\n".join(text_parts).strip()


def _load_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(paragraphs).strip()


def _load_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    text_runs: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                text_runs.append(text)
    return "\n".join(text_runs).strip()
